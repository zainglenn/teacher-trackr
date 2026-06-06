from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x1A, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xC9, 0xA0, 0x2C)
LIGHT  = RGBColor(0xF4, 0xF4, 0xF8)
MUTED  = RGBColor(0x88, 0x88, 0x99)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

SCREENSHOTS = "/Users/zainglenn/Documents/projects/education/curriculum-tracker/appScreenshots"

# ── Helpers ──────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_label(slide, text, l, t, w, h,
              size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
              wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or NAVY
    return txb

def section_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    add_label(slide, title,
              Inches(0.55), Inches(0.2), Inches(11.5), Inches(0.72),
              size=28, bold=True, color=WHITE)
    if subtitle:
        add_label(slide, subtitle,
                  Inches(0.55), Inches(0.88), Inches(11.5), Inches(0.38),
                  size=13, color=GOLD)

def footer(slide):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=NAVY)
    add_label(slide, "Zain Glenn  ·  HOD Interview 2026–2027",
              Inches(0.4), H - Inches(0.34), Inches(6.5), Inches(0.3),
              size=9, color=MUTED)
    add_label(slide, "Dubai Schools Al Khawaneej",
              Inches(7), H - Inches(0.34), Inches(6), Inches(0.3),
              size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, h, size=12, color=None, bullet="→", gap=5):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(gap)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color or NAVY
    return txb

# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(3.0), Inches(0.08), Inches(1.9), fill=GOLD)

add_label(sl, "HOD Action Plan",
          Inches(0.55), Inches(2.85), Inches(12), Inches(0.9),
          size=44, bold=True, color=WHITE)
add_label(sl, "English Middle & High School  ·  Grades 5–9",
          Inches(0.55), Inches(3.7), Inches(12), Inches(0.55),
          size=22, color=GOLD)
add_label(sl, "Dubai Schools Al Khawaneej  ·  Academic Year 2026–2027",
          Inches(0.55), Inches(4.2), Inches(12), Inches(0.4),
          size=15, color=RGBColor(0xAA, 0xAA, 0xBB))
add_label(sl, "Zain Glenn",
          Inches(0.55), Inches(5.05), Inches(6), Inches(0.4),
          size=14, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — Vision
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
section_header(sl, "Vision", "What I am here to build")

txb = sl.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.3), Inches(3.0))
tf  = txb.text_frame
tf.word_wrap = True
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = (
    "“A cohesive, standards-driven English department\n"
    "where every student is challenged by a well-sequenced curriculum,\n"
    "every teacher is growing professionally,\n"
    "and every decision is grounded in evidence — not assumption.”"
)
run.font.size   = Pt(21)
run.font.italic = True
run.font.color.rgb = NAVY

principles = [
    ("Listen before leading",  "Understand the department before changing it."),
    ("Audit before acting",    "Know the gaps before designing solutions."),
    ("Systems over heroics",   "Make excellent practice the default, not the exception."),
]
for i, (title, desc) in enumerate(principles):
    lx = Inches(0.5) + i * Inches(4.27)
    add_rect(sl, lx, Inches(5.05), Inches(3.8), Inches(1.9), fill=LIGHT)
    add_label(sl, title, lx + Inches(0.2), Inches(5.18), Inches(3.5), Inches(0.38),
              size=13, bold=True)
    add_label(sl, desc,  lx + Inches(0.2), Inches(5.58), Inches(3.5), Inches(1.1),
              size=11, color=MUTED, wrap=True)

footer(sl)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — Strategic Priorities
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
section_header(sl, "Four Strategic Priorities", "Academic Year 2026–2027")

priorities = [
    ("01", "Curriculum\nCoherence",
     "Every class across all five year groups has a complete, standards-aligned unit plan reviewed and approved before teaching begins."),
    ("02", "Differentiated\nTeacher Development",
     "Coaching calibrated to where each teacher is — not a one-size programme. Experienced teachers get leadership; developing teachers get structure."),
    ("03", "Data-Informed\nDecisions",
     "Standards coverage, attainment, and plan quality tracked across all year groups throughout the year. Evidence — not end-of-year reports — drives decisions."),
    ("04", "Grade 9\nTransition Readiness",
     "Students leaving Grade 9 are prepared for upper secondary English. Planning and assessment at this level must reflect that expectation explicitly."),
]

cw = Inches(2.95)
for i, (num, title, desc) in enumerate(priorities):
    lx = Inches(0.45) + i * Inches(3.15)
    add_rect(sl, lx, Inches(1.55), cw, Inches(4.85), fill=LIGHT)
    add_rect(sl, lx, Inches(1.55), cw, Inches(0.06), fill=GOLD)
    add_label(sl, num,   lx + Inches(0.18), Inches(1.7),  Inches(0.8), Inches(0.45),
              size=28, bold=True, color=GOLD)
    add_label(sl, title, lx + Inches(0.18), Inches(2.15), cw - Inches(0.3), Inches(0.85),
              size=14, bold=True, wrap=True)
    add_label(sl, desc,  lx + Inches(0.18), Inches(3.05), cw - Inches(0.3), Inches(3.0),
              size=11, color=RGBColor(0x44, 0x44, 0x55), wrap=True)

footer(sl)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — Term 1: Establish
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
section_header(sl, "Term 1 — Establish", "24 Aug – 6 Dec 2026  ·  13 weeks  ·  2 units per teacher")

cols = [
    ("Weeks 1–2\nFoundation", [
        "Meet every teacher individually — listen before setting direction",
        "Audit previous year’s plans: standards coverage, essential question quality, assessment alignment",
        "Begin informal classroom visits before formal observations",
    ]),
    ("Rolling\nUnit Reviews", [
        "Teachers submit each unit plan before teaching it",
        "HOD reviews: essential question, standard-to-task alignment, sequencing",
        "Approve, revise, or reject — every decision written and documented",
        "No unit taught without an approved plan",
    ]),
    ("Observation\nRound 1", [
        "Formal observation of every teacher — at least once",
        "Focus: does the classroom match the approved plan?",
        "Post-observation conversation within 48 hours — specific, not general",
    ]),
    ("End of Term\nReview", [
        "Standards coverage check across all five year groups",
        "Flag gaps to individual teachers in coaching conversations",
        "One-page summary to SLT: coverage position and PD focus for T2",
    ]),
]

cw = Inches(2.95)
for i, (heading, blist) in enumerate(cols):
    lx = Inches(0.45) + i * Inches(3.17)
    add_rect(sl, lx, Inches(1.55), cw, Inches(5.5), fill=LIGHT)
    add_rect(sl, lx, Inches(1.55), cw, Inches(0.05), fill=NAVY)
    add_label(sl, heading, lx + Inches(0.15), Inches(1.65), cw - Inches(0.25), Inches(0.72),
              size=12, bold=True, wrap=True)
    bullets(sl, blist, lx + Inches(0.15), Inches(2.45), cw - Inches(0.25), Inches(4.3),
            size=11, bullet="·", gap=5)

footer(sl)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — Terms 2 & 3
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
section_header(sl, "Terms 2 & 3 — Develop & Consolidate",
               "Term 2: 4 Jan – 14 Mar 2027  ·  Term 3: 29 Mar – 2 Jul 2027  ·  2 units per teacher per term")

# T2
add_rect(sl, Inches(0.4),  Inches(1.55), Inches(5.95), Inches(5.5), fill=LIGHT)
add_rect(sl, Inches(0.4),  Inches(1.55), Inches(5.95), Inches(0.4), fill=GREEN)
add_label(sl, "Term 2 — Develop",
          Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.32),
          size=13, bold=True, color=WHITE)

t2 = [
    "Rolling unit reviews continue — each unit approved before teaching",
    "Gather attainment data for standards taught in T1",
    "Identify patterns across classes in the same year group",
    "PLT session: assessment design — how does the task generate evidence?",
    "Observation Round 2 — compare against T1 baseline",
    "Mid-year standards coverage report to SLT",
]
bullets(sl, t2, Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.7),
        size=11.5, bullet="→", gap=6)

# T3
add_rect(sl, Inches(6.95), Inches(1.55), Inches(5.95), Inches(5.5), fill=LIGHT)
add_rect(sl, Inches(6.95), Inches(1.55), Inches(5.95), Inches(0.4), fill=AMBER)
add_label(sl, "Term 3 — Consolidate",
          Inches(7.15), Inches(1.6), Inches(5.5), Inches(0.32),
          size=13, bold=True, color=WHITE)

t3 = [
    "Rolling unit reviews continue through to final units",
    "Standards completion push — every year group targets ≥85% coverage",
    "Grade 9 transition review: analytical writing, reading, inquiry readiness",
    "Individual development reviews with every teacher",
    "Curriculum audit for 2027–2028: which standards were consistently undertaught?",
    "End-of-year report to SLT: outcomes, gaps, recommendations",
]
bullets(sl, t3, Inches(7.15), Inches(2.1), Inches(5.5), Inches(4.7),
        size=11.5, bullet="→", gap=6)

footer(sl)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — Curriculum Technology
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
section_header(sl, "Curriculum Technology",
               "Built for DSK — live at curriculum-tracker-five.vercel.app")

tx_w = Inches(5.6)
add_label(sl, "A web-based curriculum tracker I built specifically for this school context.",
          Inches(0.45), Inches(1.65), tx_w, Inches(0.5),
          size=12, bold=True, wrap=True)

features = [
    "Real-time standards coverage heatmap across all classes and year groups",
    "Unit plan submission & HOD approval workflow — every decision written and time-stamped",
    "HOD feedback recorded against each unit plan — approve, revise, or reject",
    "AI-assisted standard suggestion and gap-filling for teachers",
    "Live with Grade 6 ELA (41 standards) — ready to scale to all five year groups",
]
bullets(sl, features, Inches(0.45), Inches(2.25), tx_w, Inches(3.5),
        size=11, bullet="→", gap=7)

add_label(sl,
          "Building it required designing every HOD workflow in this plan—"
          "the review cycle, the feedback loop, the coverage audit—"
          "and making each one concrete and repeatable.",
          Inches(0.45), Inches(5.7), tx_w, Inches(0.85),
          size=10.5, italic=True, color=MUTED, wrap=True)

# Screenshots
img_l = Inches(6.35)
img_w = Inches(6.5)

try:
    sl.shapes.add_picture(f"{SCREENSHOTS}/saved-standards.jpeg",
                          img_l, Inches(1.55), img_w, Inches(2.6))
    add_label(sl, "NYSED Grade 6 ELA — 41 standards mapped by strand",
              img_l, Inches(4.2), img_w, Inches(0.28),
              size=8.5, color=MUTED, italic=True)
except Exception as e:
    print(f"img1 error: {e}")

try:
    sl.shapes.add_picture(f"{SCREENSHOTS}/ai-preview-1280.jpeg",
                          img_l, Inches(4.55), img_w, Inches(2.2))
    add_label(sl, "AI-generated standard set — reviewed and confirmed before saving",
              img_l, Inches(6.78), img_w, Inches(0.28),
              size=8.5, color=MUTED, italic=True)
except Exception as e:
    print(f"img2 error: {e}")

footer(sl)

# ── Save ─────────────────────────────────────────────────
out = "/Users/zainglenn/Documents/projects/education/curriculum-tracker/.design/HOD-ACTION-PLAN-2026-27.pptx"
prs.save(out)
print(f"Saved → {out}")
